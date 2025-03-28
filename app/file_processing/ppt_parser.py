from typing import Dict
from pptx import Presentation
from .base_parser import BaseParser

class PPTParser(BaseParser):
    def parse(self, file_path: str) -> Dict:
        try:
            prs = Presentation(file_path)
            return {
                'text': '\n'.join([slide.shapes.title.text for slide in prs.slides if slide.shapes.title]),
                'metadata': {
                    'slides': [{
                        'layout': slide.slide_layout.name,
                        'notes': slide.notes_slide.notes_text_frame.text
                    } for slide in prs.slides]
                }
            }
        except Exception as e:
            self.logger.error(f"PPT解析失败: {str(e)}")
            raise